"""Generate editable PPTX from slide outline for the NeurIPS 2026 paper."""
from pathlib import Path
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---- NeurIPS color scheme ----
COL_PRIMARY = RGBColor(0x8B, 0x5C, 0xF6)
COL_ACCENT  = RGBColor(0x25, 0x63, 0xEB)
COL_TEXT    = RGBColor(0x1E, 0x1E, 0x1E)
COL_GRAY    = RGBColor(0x64, 0x64, 0x64)

SLIDES_DIR = Path(__file__).parent
FIG_DIR = SLIDES_DIR / "figures"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def ensure_png(stem):
    """Rasterize PDF to PNG at high resolution if not already."""
    png = FIG_DIR / f"{stem}.png"
    pdf = FIG_DIR / f"{stem}.pdf"
    if png.exists() or not pdf.exists():
        return png if png.exists() else None
    try:
        subprocess.run(
            ["pdftoppm", "-png", "-r", "220", str(pdf), str(FIG_DIR / stem)],
            check=True, capture_output=True,
        )
        generated = FIG_DIR / f"{stem}-1.png"
        if generated.exists():
            generated.rename(png)
            return png
    except Exception:
        return None
    return None


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_frame_number(slide, n, total):
    tb = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(1.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {total}"
    r.font.size = Pt(10)
    r.font.color.rgb = COL_GRAY


def add_title(slide, text, size=32):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = COL_PRIMARY


def add_underline_bar(slide, y=Inches(1.2), width=Inches(2.2)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, width, Emu(38100))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COL_PRIMARY
    bar.line.fill.background()


def add_bullets(slide, bullets, top=Inches(1.5), size=22):
    tb = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = COL_TEXT


def add_text(slide, text, left, top, width, height, size=20, bold=False, color=COL_TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_image(slide, png, left, top, width=None, height=None):
    if png is None or not png.exists():
        return
    if width:
        slide.shapes.add_picture(str(png), left, top, width=width)
    elif height:
        slide.shapes.add_picture(str(png), left, top, height=height)
    else:
        slide.shapes.add_picture(str(png), left, top)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


TOTAL = 15

# -------- 1. Title --------
s = add_slide()
add_text(s, "Rethinking LLM Internal-State Probing",
         Inches(0.6), Inches(2.4), Inches(12.1), Inches(1.0),
         size=44, bold=True, color=COL_PRIMARY, align=PP_ALIGN.CENTER)
add_text(s, "From Fragmented Signals to Extensible Fusion",
         Inches(0.6), Inches(3.4), Inches(12.1), Inches(0.8),
         size=28, color=COL_ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "Junyi Chen",
         Inches(0.6), Inches(4.8), Inches(12.1), Inches(0.5),
         size=22, align=PP_ALIGN.CENTER)
add_text(s, "University of Auckland    ·    NeurIPS 2026",
         Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.5),
         size=18, color=COL_GRAY, align=PP_ALIGN.CENTER)
set_notes(s, "Thank the chair, introduce self and paper title. Keep brief.")

# -------- 2. Why Probe --------
s = add_slide()
add_title(s, "Why Probe LLMs?")
add_underline_bar(s)
add_bullets(s, [
    "Many deployment-critical properties are latent: hallucination, factuality, calibration, tool invocation",
    "Probing reads an internal signal from a frozen LLM and fits a lightweight classifier",
    "No finetuning, cheap to scale: a default tool in modern LLM analysis",
])
add_frame_number(s, 2, TOTAL)
set_notes(s, "Frame probing as the go-to technique. The properties we care about are inside the model.")

# -------- 3. Hero figure (only visual for the reframing) --------
s = add_slide()
add_title(s, "From Fragmented Probes to an Extensible Scaffold")
add_underline_bar(s, width=Inches(4.8))
hero_png = ensure_png("fig0_hero")
add_image(s, hero_png, Inches(2.0), Inches(1.5), width=Inches(9.3))
add_text(s,
    "Top: current paradigm, probes read in isolation.    Bottom: ideal paradigm, probes plug into a shared scaffold.",
    Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4),
    size=14, color=COL_GRAY, align=PP_ALIGN.CENTER)
add_frame_number(s, 3, TOTAL)
set_notes(s, "Thesis slide. Walk top (fragmentation) then bottom (plug-in scaffold). The single visual the talk hinges on.")

# -------- 4. RQ1 setup --------
s = add_slide()
add_title(s, "RQ1: Does Any Probe Generalize?")
add_underline_bar(s, width=Inches(4.2))
add_bullets(s, [
    "Benchmark: 7 representative probes × 5 classification datasets × 3 open LLMs",
    "Tasks span hallucination, factuality, math, tool routing",
    "Each probe uses its original pipeline and hyperparameters",
], top=Inches(1.5), size=22)
bh_png = ensure_png("fig_baseline_heatmap")
add_image(s, bh_png, Inches(1.3), Inches(4.1), width=Inches(10.7))
add_frame_number(s, 4, TOTAL)
set_notes(s, "Empirical setup. Hold on the heatmap briefly; punchline is on slide 5.")

# -------- 5. RQ1 answer --------
s = add_slide()
add_title(s, "No Single Probe Dominates")
add_underline_bar(s, width=Inches(3.7))
add_image(s, bh_png, Inches(1.3), Inches(1.6), width=Inches(10.7))
add_bullets(s, [
    "Column-winning probe shifts across tasks and models",
    "The ranking of probes changes in almost every cell",
], top=Inches(5.7), size=20)
add_frame_number(s, 5, TOTAL)
set_notes(s, "Deliver the punchline: fragmentation is real. Point to cells where the winner flips.")

# -------- 6. RQ2 setup --------
s = add_slide()
add_title(s, "RQ2: Is Probe Disagreement Exploitable?")
add_underline_bar(s, width=Inches(5.6))
add_bullets(s, [
    "Per-example oracle: for each test example, pick the probe that gets it right",
    "Headroom = oracle AUROC − best-single-probe AUROC",
    "Small headroom: probes agree. Large headroom: probes make genuinely different errors.",
], top=Inches(1.5), size=22)
or_png = ensure_png("fig4_oracle")
add_image(s, or_png, Inches(2.5), Inches(4.4), width=Inches(8.3))
add_frame_number(s, 6, TOTAL)
set_notes(s, "Define the oracle. It measures whether fusion even has a ceiling worth chasing.")

# -------- 7. Oracle headroom --------
s = add_slide()
add_title(s, "Probes Make Genuinely Different Errors")
add_underline_bar(s, width=Inches(5.1))
add_image(s, or_png, Inches(2.5), Inches(1.6), width=Inches(8.3))
add_bullets(s, [
    "Oracle headroom exceeds 10 pp on every dataset",
    "Simple stacking already captures part of it: the complementarity is exploitable in practice",
], top=Inches(5.7), size=20)
add_frame_number(s, 7, TOTAL)
set_notes(s, "Green-vs-blue gap is the ceiling. Orange bar shows simple stacking already benefits.")

# -------- 8. RQ3 structure --------
s = add_slide()
add_title(s, "RQ3: Where Does the Diversity Come From?")
add_underline_bar(s, width=Inches(6.0))
cl_png = ensure_png("fig3_clustering")
add_image(s, cl_png, Inches(1.0), Inches(1.6), width=Inches(11.3))
add_bullets(s, [
    "Probes cluster by signal family: hidden-state, attention, generation-side",
    "Cross-family correlation is low: diversity is structural, not random noise",
], top=Inches(6.1), size=18)
add_frame_number(s, 8, TOTAL)
set_notes(s, "Spearman between probes. Clusters are clean and consistent across models.")

# -------- 9. Competence --------
s = add_slide()
add_title(s, "But Orthogonality Is Not Enough")
add_underline_bar(s, width=Inches(4.8))
cp_png = ensure_png("fig_competence")
add_image(s, cp_png, Inches(1.0), Inches(1.6), width=Inches(11.3))
add_bullets(s, [
    "Most-orthogonal probes (generation-side) contribute the least because they are weakest in isolation",
    "Effective fusion needs diversity AND individual probe competence",
], top=Inches(6.1), size=18)
add_frame_number(s, 9, TOTAL)
set_notes(s, "Counterintuitive: more orthogonal is not more useful. Competence is the missing variable.")

# -------- 10. Design targets --------
s = add_slide()
add_title(s, "Design Targets for the Scaffold")
add_underline_bar(s, width=Inches(4.7))
add_bullets(s, [
    "(i)  Decision-level fusion  →  keep each probe's prediction separate, preserve cross-family orthogonality",
    "(ii) Learned meta-combiner  →  per-task weighting, no manual probe selection",
    "(iii) Plug-in extensibility  →  adding a new probe requires retraining only the combiner",
], top=Inches(1.9), size=24)
add_frame_number(s, 10, TOTAL)
set_notes(s, "The three design targets fall out of the three findings.")

# -------- 11. Two-stage pipeline --------
def stage_box(slide, left, top, w, h, title, body, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF3, 0xFF) if color == COL_PRIMARY else RGBColor(0xEF, 0xF6, 0xFF)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.25)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = body
    r2.font.size = Pt(17); r2.font.color.rgb = COL_TEXT

s = add_slide()
add_title(s, "The Scaffold: Two Stages")
add_underline_bar(s, width=Inches(3.6))
stage_box(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(4.0),
          "Stage 1 — Per-expert prediction",
          "Each probe wrapped in an identical predictor:\n• StandardScaler + PCA\n• Logistic regression (linear)\n• Gradient-boosted trees (nonlinear)\n\n→ 5-fold out-of-fold probabilities per probe",
          COL_PRIMARY)
stage_box(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.0),
          "Stage 2 — Meta-classification",
          "Concatenate all per-expert probabilities\ninto a meta-feature vector\n\nFit L2-regularized logistic regression\nwith cross-validated strength\n\n→ Final class prediction",
          COL_ACCENT)
add_text(s,
    "Adding a new probe: train its expert, refit the combiner. Existing experts are untouched.",
    Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6),
    size=18, color=COL_GRAY, align=PP_ALIGN.CENTER)
add_frame_number(s, 11, TOTAL)
set_notes(s, "Walk through Stage 1 then Stage 2. Emphasize decision-level fusion and plug-in property.")

# -------- 12. Main results --------
s = add_slide()
add_title(s, "Main Results")
add_underline_bar(s, width=Inches(2.0))
add_text(s, "Ranks first on all 15 dataset-model pairs",
         Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.7),
         size=28, bold=True, color=COL_PRIMARY, align=PP_ALIGN.CENTER)

def big_num(slide, left, num, label):
    add_text(slide, num, left, Inches(3.0), Inches(4.0), Inches(1.1),
             size=56, bold=True, color=COL_ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, Inches(4.2), Inches(4.0), Inches(0.6),
             size=18, color=COL_TEXT, align=PP_ALIGN.CENTER)

big_num(s, Inches(0.3),  "+1.7%",  "Median gain over best single probe")
big_num(s, Inches(4.67), "+5.6%",  "Largest single gain  (Qwen / When2Call)")
big_num(s, Inches(9.03), "15 / 15", "Dataset-model pairs improved")

add_text(s,
    "No task-specific tuning. Same pipeline, same hyperparameters across datasets and models.",
    Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.6),
    size=20, color=COL_GRAY, align=PP_ALIGN.CENTER)
add_frame_number(s, 12, TOTAL)
set_notes(s, "Headline result. 15/15 is every dataset-model combination tested.")

# -------- 13. Ablation --------
s = add_slide()
add_title(s, "Both Design Choices Are Load-Bearing")
add_underline_bar(s, width=Inches(4.8))
add_bullets(s, [
    "Meta-combiner:  replacing the learned combiner with Prediction Averaging loses up to 6.1 pp",
    "Simple LR Stacking recovers most of it but still fails on When2Call (heterogeneous probe behavior)",
    "Expert library:  LR-only or GBT-only loses several pp on heterogeneous tasks; LR + GBT covers both",
    "Neural meta-learners (≈ 500K params) degrade by 2–9 pp: standard overfitting at these training sizes",
], top=Inches(1.7), size=22)
add_frame_number(s, 13, TOTAL)
set_notes(s, "Naive fusion hurts. Library diversity is the lift on heterogeneous tasks.")

# -------- 14. Takeaway --------
s = add_slide()
add_title(s, "Takeaway")
add_underline_bar(s, width=Inches(1.7))
add_text(s, "Stop searching for a universal probe.",
         Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.8),
         size=32, bold=True, color=COL_PRIMARY, align=PP_ALIGN.CENTER)
add_text(s, "Maintain a scaffold that absorbs each new signal as a plug-in expert.",
         Inches(0.6), Inches(3.2), Inches(12.1), Inches(0.8),
         size=28, color=COL_ACCENT, align=PP_ALIGN.CENTER)

def recap(slide, left, head, body):
    add_text(slide, head, left, Inches(5.0), Inches(4.0), Inches(0.6),
             size=22, bold=True, color=COL_TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, body, left, Inches(5.6), Inches(4.0), Inches(1.3),
             size=16, color=COL_GRAY, align=PP_ALIGN.CENTER)

recap(s, Inches(0.3),  "Diagnostic",   "No single probe generalizes;\noracle headroom is large")
recap(s, Inches(4.67), "Structural",   "Disagreement clusters\nby signal family")
recap(s, Inches(9.03), "Constructive", "Two-stage scaffold wins\non all 15 pairs")
add_frame_number(s, 14, TOTAL)
set_notes(s, "End with the one-liner. Recap three contributions in one line each.")

# -------- 15. Thank you --------
s = add_slide()
add_text(s, "Thank You", Inches(0.6), Inches(2.6), Inches(12.1), Inches(1.2),
         size=60, bold=True, color=COL_PRIMARY, align=PP_ALIGN.CENTER)
add_text(s, "Questions?", Inches(0.6), Inches(3.9), Inches(12.1), Inches(0.7),
         size=32, color=COL_ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "Paper + code available", Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.5),
         size=20, align=PP_ALIGN.CENTER)
add_text(s, "(QR code placeholder — replace with real URL)",
         Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.5),
         size=14, color=COL_GRAY, align=PP_ALIGN.CENTER)
set_notes(s, "Pause and invite questions. Swap placeholder for the real QR before presenting.")

out = SLIDES_DIR / "presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
